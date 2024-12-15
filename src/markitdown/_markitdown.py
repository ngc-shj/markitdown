# type: ignore
import base64
import binascii
from collections import defaultdict
import copy
from dataclasses import dataclass
import hashlib
import html
import json
import mimetypes
import os
import yaml
import re
import shutil
import sqlparse
import subprocess
import sys
import tempfile
import traceback
from typing import Any, Dict, List, Optional, Union
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse
import xml.etree.ElementTree as ET

import mammoth
import markdownify
import pandas as pd
import pdfminer
import pdfminer.high_level
import pptx

# File-format detection
import puremagic
import requests
from bs4 import BeautifulSoup

# Optional Transcription support
try:
    import pydub
    import speech_recognition as sr

    IS_AUDIO_TRANSCRIPTION_CAPABLE = True
except ModuleNotFoundError:
    pass

# Optional YouTube transcription support
try:
    from youtube_transcript_api import YouTubeTranscriptApi

    IS_YOUTUBE_TRANSCRIPT_CAPABLE = True
except ModuleNotFoundError:
    pass


class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        # Explicitly cast options to the expected type if necessary
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

        return super().convert_hn(n, el, text, convert_as_inline)  # type: ignore

    def convert_a(self, el: Any, text: str, convert_as_inline: bool):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)  # type: ignore
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        # Escape URIs and skip non-http or file schemes
        if href:
            try:
                parsed_url = urlparse(href)  # type: ignore
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:  # type: ignore
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))  # type: ignore
            except ValueError:  # It's not clear if this ever gets thrown
                return "%s%s%s" % (prefix, text, suffix)

        # For the replacement see #29: text nodes underscores are escaped
        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            # Shortcut syntax
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return (
            "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix)
            if href
            else text
        )

    def convert_img(self, el: Any, text: str, convert_as_inline: bool) -> str:
        """Same as usual converter, but removes data URIs"""

        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if (
            convert_as_inline
            and el.parent.name not in self.options["keep_inline_images_in"]
        ):
            return alt

        # Remove dataURIs
        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)  # type: ignore


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title: Union[str, None] = title
        self.text_content: str = text_content


class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Guess the content type from any file extension that might be around
        content_type, _ = mimetypes.guess_type(
            "__placeholder" + kwargs.get("file_extension", "")
        )

        # Only accept text files
        if content_type is None:
            return None
        elif "text/" not in content_type.lower():
            return None

        text_content = ""
        with open(local_path, "rt", encoding="utf-8") as fh:
            text_content = fh.read()
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content: str) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        assert isinstance(webpage_text, str)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not Wikipedia
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        main_title = None if soup.title is None else soup.title.string

        if body_elm:
            # What's the title
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string  # type: ignore
                assert isinstance(main_title, str)

            # Convert the page
            webpage_text = f"# {main_title}\n\n" + _CustomMarkdownify().convert_soup(
                body_elm
            )
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=main_title,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert(
        self, local_path: str, **kwargs: Any
    ) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        assert soup.title is not None and soup.title.string is not None
        metadata: Dict[str, str] = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")  # type: ignore
                        if attrdesc:
                            metadata["description"] = str(attrdesc["content"])
                    break
        except Exception:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])  # type: ignore
        assert isinstance(title, str)

        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])  # type: ignore
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])  # type: ignore
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])  # type: ignore
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])  # type: ignore
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        if IS_YOUTUBE_TRANSCRIPT_CAPABLE:
            transcript_text = ""
            parsed_url = urlparse(url)  # type: ignore
            params = parse_qs(parsed_url.query)  # type: ignore
            if "v" in params:
                assert isinstance(params["v"][0], str)
                video_id = str(params["v"][0])
                try:
                    # Must be a single transcript.
                    transcript = YouTubeTranscriptApi.get_transcript(video_id)  # type: ignore
                    transcript_text = " ".join([part["text"] for part in transcript])  # type: ignore
                    # Alternative formatting:
                    # formatter = TextFormatter()
                    # formatter.format_transcript(transcript)
                except Exception:
                    pass
            if transcript_text:
                webpage_text += f"\n### Transcript\n{transcript_text}\n"

        title = title if title else soup.title.string
        assert isinstance(title, str)

        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
        )

    def _get(
        self,
        metadata: Dict[str, str],
        keys: List[str],
        default: Union[str, None] = None,
    ) -> Union[str, None]:
        for k in keys:
            if k in metadata:
                return metadata[k]
        return default

    def _findKey(self, json: Any, key: str) -> Union[str, None]:  # TODO: Fix json type
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None


class BingSerpConverter(DocumentConverter):
    """
    Handle Bing results pages (only the organic search results).
    NOTE: It is better to use the Bing API
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a Bing SERP
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https://www\.bing\.com/search\?q=", url):
            return None

        # Parse the query parameters
        parsed_params = parse_qs(urlparse(url).query)
        query = parsed_params.get("q", [""])[0]

        # Parse the file
        soup = None
        with open(local_path, "rt", encoding="utf-8") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Clean up some formatting
        for tptt in soup.find_all(class_="tptt"):
            if hasattr(tptt, "string") and tptt.string:
                tptt.string += " "
        for slug in soup.find_all(class_="algoSlug_icon"):
            slug.extract()

        # Parse the algorithmic results
        _markdownify = _CustomMarkdownify()
        results = list()
        for result in soup.find_all(class_="b_algo"):
            # Rewrite redirect urls
            for a in result.find_all("a", href=True):
                parsed_href = urlparse(a["href"])
                qs = parse_qs(parsed_href.query)

                # The destination is contained in the u parameter,
                # but appears to be base64 encoded, with some prefix
                if "u" in qs:
                    u = (
                        qs["u"][0][2:].strip() + "=="
                    )  # Python 3 doesn't care about extra padding

                    try:
                        # RFC 4648 / Base64URL" variant, which uses "-" and "_"
                        a["href"] = base64.b64decode(u, altchars="-_").decode("utf-8")
                    except UnicodeDecodeError:
                        pass
                    except binascii.Error:
                        pass

            # Convert to markdown
            md_result = _markdownify.convert_soup(result).strip()
            lines = [line.strip() for line in re.split(r"\n+", md_result)]
            results.append("\n".join([line for line in lines if len(line) > 0]))

        webpage_text = (
            f"## A Bing search for '{query}' found the following results:\n\n"
            + "\n\n".join(results)
        )

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(HtmlConverter):
    """
    Converts XLSX files to Markdown, with each sheet presented as a separate Markdown table.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += (
                        "\n!["
                        + (alt_text if alt_text else shape.name)
                        + "]("
                        + filename
                        + ")\n"
                    )

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += (
                        "\n" + self._convert(html_table).text_content.strip() + "\n"
                    )

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run(
                    [exiftool, "-json", local_path], capture_output=True, text=True
                ).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            try:
                transcript = self._transcribe_audio(local_path)
                md_content += "\n\n### Audio Transcript:\n" + (
                    "[No speech detected]" if transcript == "" else transcript
                )
            except Exception:
                md_content += (
                    "\n\n### Audio Transcript:\nError. Could not transcribe this audio."
                )

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        recognizer = sr.Recognizer()
        with sr.AudioFile(local_path) as source:
            audio = recognizer.record(source)
            return recognizer.recognize_google(audio).strip()


class Mp3Converter(WavConverter):
    """
    Converts MP3 files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".mp3":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        if IS_AUDIO_TRANSCRIPTION_CAPABLE:
            handle, temp_path = tempfile.mkstemp(suffix=".wav")
            os.close(handle)
            try:
                sound = pydub.AudioSegment.from_mp3(local_path)
                sound.export(temp_path, format="wav")

                _args = dict()
                _args.update(kwargs)
                _args["file_extension"] = ".wav"

                try:
                    transcript = super()._transcribe_audio(temp_path).strip()
                    md_content += "\n\n### Audio Transcript:\n" + (
                        "[No speech detected]" if transcript == "" else transcript
                    )
                except Exception:
                    md_content += "\n\n### Audio Transcript:\nError. Could not transcribe this audio."

            finally:
                os.unlink(temp_path)

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an mlm_client is configured).
    """

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        mlm_client = kwargs.get("mlm_client")
        mlm_model = kwargs.get("mlm_model")
        if mlm_client is not None and mlm_model is not None:
            md_content += (
                "\n# Description:\n"
                + self._get_mlm_description(
                    local_path,
                    extension,
                    mlm_client,
                    mlm_model,
                    prompt=kwargs.get("mlm_prompt"),
                ).strip()
                + "\n"
            )

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _get_mlm_description(self, local_path, extension, client, model, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        sys.stderr.write(f"MLM Prompt:\n{prompt}\n")

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": data_uri,
                        },
                    },
                ],
            }
        ]

        response = client.chat.completions.create(model=model, messages=messages)
        return response.choices[0].message.content


class JsonConverter(DocumentConverter):
    """Converts JSON files to a readable Markdown format"""
    
    def convert(self, local_path: str, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".json":
            return None
            
        try:
            with open(local_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            md_content = "# JSON Contents\n\n"
            md_content += self._json_to_markdown(data)
            
            return DocumentConverterResult(
                title="JSON Document",
                text_content=md_content
            )
        except json.JSONDecodeError:
            return None
            
    def _json_to_markdown(self, data, level=0):
        md = ""
        indent = "  " * level
        
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    md += f"{indent}- **{key}**:\n{self._json_to_markdown(value, level+1)}"
                else:
                    md += f"{indent}- **{key}**: {value}\n"
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, (dict, list)):
                    md += f"{indent}-\n{self._json_to_markdown(item, level+1)}"
                else:
                    md += f"{indent}- {item}\n"
        return md


class EnhancedCsvConverter(DocumentConverter):
    """Converts CSV files to Markdown tables with comprehensive data analysis"""
    
    def convert(self, local_path: str, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".csv":
            return None
            
        try:
            # Use pandas for data analysis
            import pandas as pd
            df = pd.read_csv(local_path)
            
            md_content = "# CSV Analysis\n\n"
            
            # Basic dataset information
            md_content += "## Summary Statistics\n\n"
            md_content += f"- Total Rows: {len(df)}\n"
            md_content += f"- Total Columns: {len(df.columns)}\n"
            md_content += f"- Column Names: {', '.join(df.columns)}\n\n"

            # Detailed analysis for each column
            md_content += "## Column Analysis\n\n"
            for column in df.columns:
                md_content += f"### {column}\n\n"
                
                # Handle numeric data
                if pd.api.types.is_numeric_dtype(df[column]):
                    md_content += f"- Type: Numeric\n"
                    md_content += f"- Mean: {df[column].mean():.2f}\n"
                    md_content += f"- Median: {df[column].median():.2f}\n"
                    md_content += f"- Std Dev: {df[column].std():.2f}\n"
                    md_content += f"- Min: {df[column].min()}\n"
                    md_content += f"- Max: {df[column].max()}\n"
                # Handle datetime data
                elif pd.api.types.is_datetime64_any_dtype(df[column]):
                    md_content += f"- Type: DateTime\n"
                    md_content += f"- Earliest: {df[column].min()}\n"
                    md_content += f"- Latest: {df[column].max()}\n"
                # Handle other data types (categorical, text, etc.)
                else:
                    md_content += f"- Type: {df[column].dtype}\n"
                    md_content += f"- Unique Values: {df[column].nunique()}\n"
                    # Show top 5 most frequent values
                    value_counts = df[column].value_counts().head()
                    if not value_counts.empty:
                        md_content += "- Top Values:\n"
                        for val, count in value_counts.items():
                            md_content += f"  - {val}: {count} times\n"
                
                md_content += "\n"

            # Data preview section
            preview_rows = 5  # Default number of rows to show at start and end
            
            md_content += "## Data Preview\n\n"
            md_content += "### First {} Rows\n\n".format(preview_rows)
            
            # Add table headers
            md_content += "| " + " | ".join(df.columns) + " |\n"
            md_content += "| " + " | ".join(["---"] * len(df.columns)) + " |\n"
            
            # Show first n rows
            for _, row in df.head(preview_rows).iterrows():
                md_content += "| " + " | ".join(str(val) for val in row) + " |\n"
            
            # For large datasets, also show the last few rows
            if len(df) > preview_rows * 2:
                md_content += "\n### Last {} Rows\n\n".format(preview_rows)
                
                # Repeat headers for the bottom section
                md_content += "| " + " | ".join(df.columns) + " |\n"
                md_content += "| " + " | ".join(["---"] * len(df.columns)) + " |\n"
                
                # Show last n rows
                for _, row in df.tail(preview_rows).iterrows():
                    md_content += "| " + " | ".join(str(val) for val in row) + " |\n"
            
            # Add note for truncated data
            if len(df) > preview_rows * 2:
                total_shown = preview_rows * 2
                md_content += f"\n*Note: Showing {total_shown} rows out of {len(df)} total rows*\n"
            
            return DocumentConverterResult(
                title="CSV Analysis",
                text_content=md_content
            )
        except Exception as e:
            # Log error and delegate to other converters if available
            print(f"Error in CSV conversion: {str(e)}")
            return None


class XmlConverter(DocumentConverter):
    """Converts XML files to a structured Markdown format"""
    
    def convert(self, local_path: str, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xml":
            return None
            
        try:
            tree = ET.parse(local_path)
            root = tree.getroot()
            
            md_content = "# XML Document Structure\n\n"
            md_content += self._element_to_markdown(root)
            
            return DocumentConverterResult(
                title="XML Document",
                text_content=md_content
            )
        except ET.ParseError:
            return None
            
    def _element_to_markdown(self, element, level=0):
        indent = "  " * level
        md = f"{indent}- Element: **{element.tag}**\n"
        
        # Add attributes if any
        if element.attrib:
            md += f"{indent}  Attributes:\n"
            for key, value in element.attrib.items():
                md += f"{indent}    - {key}: {value}\n"
        
        # Add text content if any
        if element.text and element.text.strip():
            md += f"{indent}  Text: {element.text.strip()}\n"
        
        # Process child elements
        for child in element:
            md += self._element_to_markdown(child, level + 1)
            
        return md


class YamlConverter(DocumentConverter):
    """Converts YAML files to a readable Markdown format"""
    
    def convert(self, local_path: str, **kwargs) -> Union[None, DocumentConverterResult]:
        # Check if this is a YAML file
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".yaml", ".yml"]:
            return None
            
        try:
            with open(local_path, 'r', encoding='utf-8') as f:
                # Load all documents in the YAML file
                docs = list(yaml.safe_load_all(f))
            
            # Start building markdown content
            md_content = ""
            
            # Handle multiple documents in the YAML file
            if len(docs) > 1:
                for i, doc in enumerate(docs, 1):
                    md_content += f"# Document {i}\n\n"
                    md_content += self._detect_and_format_yaml(doc)
                    md_content += "\n---\n\n"  # Document separator
            else:
                doc = docs[0]
                md_content = self._detect_and_format_yaml(doc)
            
            return DocumentConverterResult(
                title="YAML Document",
                text_content=md_content
            )
        except yaml.YAMLError:
            return None
    
    def _detect_and_format_yaml(self, data: Any) -> str:
        """Detect common YAML formats and apply appropriate formatting"""
        
        # Detect Kubernetes manifest
        if isinstance(data, dict) and all(k in data for k in ['apiVersion', 'kind', 'metadata']):
            return self._format_kubernetes_manifest(data)
        
        # Detect Docker Compose
        elif isinstance(data, dict) and 'services' in data:
            return self._format_docker_compose(data)
        
        # Detect Ansible playbook
        elif isinstance(data, list) and all(isinstance(item, dict) and 'name' in item for item in data):
            return self._format_ansible_playbook(data)
        
        # Default formatting for general YAML
        else:
            return self._yaml_to_markdown(data)
    
    def _format_kubernetes_manifest(self, data: Dict) -> str:
        """Special formatting for Kubernetes manifests"""
        md = f"# Kubernetes {data['kind']}\n\n"
        md += f"- **API Version:** {data['apiVersion']}\n"
        md += f"- **Kind:** {data['kind']}\n"
        md += f"- **Name:** {data['metadata'].get('name', 'N/A')}\n"
        md += f"- **Namespace:** {data['metadata'].get('namespace', 'default')}\n\n"
        
        # Add other sections
        for key, value in data.items():
            if key not in ['apiVersion', 'kind', 'metadata']:
                md += f"## {key}\n\n"
                md += self._yaml_to_markdown(value, level=0)
                md += "\n"
        
        return md
    
    def _format_docker_compose(self, data: Dict) -> str:
        """Special formatting for Docker Compose files"""
        md = "# Docker Compose Configuration\n\n"
        
        if 'version' in data:
            md += f"**Version:** {data['version']}\n\n"
        
        md += "## Services\n\n"
        for service_name, config in data['services'].items():
            md += f"### {service_name}\n\n"
            md += self._yaml_to_markdown(config, level=0)
            md += "\n"
        
        # Add other top-level sections
        for key, value in data.items():
            if key not in ['version', 'services']:
                md += f"## {key}\n\n"
                md += self._yaml_to_markdown(value, level=0)
                md += "\n"
        
        return md
    
    def _format_ansible_playbook(self, data: List) -> str:
        """Special formatting for Ansible playbooks"""
        md = "# Ansible Playbook\n\n"
        
        for task in data:
            md += f"## Task: {task['name']}\n\n"
            task_copy = task.copy()
            task_copy.pop('name', None)  # Remove name as it's already in the header
            md += self._yaml_to_markdown(task_copy, level=0)
            md += "\n"
        
        return md
    
    def _yaml_to_markdown(self, data: Any, level: int = 0) -> str:
        """Convert YAML data to Markdown with proper formatting"""
        md = ""
        indent = "  " * level
        
        if isinstance(data, dict):
            for key, value in data.items():
                if isinstance(value, (dict, list)):
                    md += f"{indent}- **{key}**:\n{self._yaml_to_markdown(value, level+1)}"
                else:
                    md += f"{indent}- **{key}**: {value}\n"
        elif isinstance(data, list):
            for item in data:
                if isinstance(item, (dict, list)):
                    md += f"{indent}-\n{self._yaml_to_markdown(item, level+1)}"
                else:
                    md += f"{indent}- {item}\n"
        return md


@dataclass
class TableStructure:
    """Represents a database table structure"""
    name: str
    columns: List[Dict[str, str]]
    primary_keys: List[str]
    foreign_keys: List[Dict[str, List[str]]]


class SqlConverter(DocumentConverter):
    """Converts SQL files to readable Markdown format."""

    def convert(self, local_path: str, **kwargs: Any) -> Union[None, DocumentConverterResult]:
        """Convert SQL file to Markdown."""
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".sql":
            return None

        try:
            with open(local_path, "rt", encoding="utf-8") as sql_file:
                content = sql_file.read()
                return self._convert_sql_to_markdown(content)
        except Exception as e:
            return DocumentConverterResult(
                title="SQL File",
                text_content=f"Error converting SQL file: {str(e)}"
            )

    def _convert_sql_to_markdown(self, sql_content: str) -> DocumentConverterResult:
        """Convert SQL content to Markdown format."""
        markdown_content = ["# SQL File Contents\n"]
        
        statements = sqlparse.parse(sql_content)
        for i, statement in enumerate(statements, 1):
            if not statement.tokens:
                continue

            statement_type = statement.get_type() or "Statement"
            markdown_content.append(f"## {statement_type} {i}\n")

            # Format SQL
            formatted_sql = self._format_statement(statement)
            markdown_content.append(f"```sql\n{formatted_sql}\n```\n")

            # Add table documentation for CREATE statements
            if statement.get_type() == "CREATE":
                table_structure = self._parse_create_table(statement)
                if table_structure:
                    doc = self._format_table_structure(table_structure)
                    markdown_content.append(f"{doc}\n")

        return DocumentConverterResult(
            title="SQL Documentation",
            text_content='\n'.join(markdown_content).strip()
        )

    def _format_statement(self, statement: sqlparse.sql.Statement) -> str:
        """Format SQL statement with proper indentation."""
        sql_str = str(statement)
        comment = ""

        # Extract comment if present
        if '--' in sql_str:
            parts = sql_str.split('--', 1)
            sql_str = parts[0].strip()
            comment = f"-- {parts[1].strip()}\n"

        if sql_str.upper().startswith('CREATE TABLE'):
            return comment + self._format_create_table(sql_str)
        elif sql_str.upper().startswith('SELECT'):
            return comment + self._format_select(sql_str)
        else:
            return comment + sqlparse.format(sql_str, reindent=True, 
                                          keyword_case='upper', identifier_case='lower')

    def _format_create_table(self, sql: str) -> str:
        """Format CREATE TABLE statement."""
        # Normalize whitespace
        sql = ' '.join(sql.split())
        
        # Extract table name and body
        match = re.match(
            r"CREATE\s+TABLE\s+(\w+)\s*\((.*)\)",
            sql,
            re.IGNORECASE | re.DOTALL
        )
        if not match:
            return sql

        table_name = match.group(1)
        definitions = self._split_definitions(match.group(2))
        
        # Separate columns and constraints
        columns = []
        constraints = []
        for d in definitions:
            if d.upper().startswith(('PRIMARY KEY', 'FOREIGN KEY', 'CONSTRAINT')):
                constraints.append(d)
            else:
                columns.append(d)

        # Build formatted output
        lines = [f"CREATE TABLE {table_name} ("]
        
        # Add columns
        for col in columns:
            lines.append(f"    {col},")
        
        # Add constraints
        if constraints:
            if columns:  # Add blank line if we have both columns and constraints
                lines[-1] = lines[-1]  # Keep the comma
                lines.append("")
            
            for i, constraint in enumerate(constraints):
                comma = ',' if i < len(constraints) - 1 else ''
                lines.append(f"    {constraint}{comma}")
        
        lines.append(");")
        
        return '\n'.join(lines)

    def _split_definitions(self, body: str) -> List[str]:
        """Split column/constraint definitions handling nested parentheses."""
        definitions = []
        current = []
        depth = 0
        
        for char in body + ',':
            if char == '(':
                depth += 1
            elif char == ')':
                depth -= 1
                
            if char == ',' and depth == 0:
                if current:
                    definitions.append(''.join(current).strip())
                current = []
            else:
                current.append(char)
                
        return definitions

    def _format_select(self, sql: str) -> str:
        """Format SELECT statement."""
        return sqlparse.format(
            sql,
            reindent=True,
            keyword_case='upper',
            identifier_case='lower'
        )

    def _parse_create_table(self, statement: sqlparse.sql.Statement) -> Optional[TableStructure]:
        """Extract table structure from CREATE TABLE statement."""
        sql_str = str(statement)
        
        # Parse table name and body
        match = re.search(
            r"CREATE\s+TABLE\s+([^\s(]+)\s*\((.*)\)",
            sql_str,
            re.IGNORECASE | re.DOTALL
        )
        if not match:
            return None

        table_name = match.group(1).strip('`"[]')
        definitions = self._split_definitions(match.group(2))

        # Parse definitions
        columns = []
        primary_keys = []
        foreign_keys = []

        for defn in definitions:
            defn = defn.strip()
            
            # Skip empty definitions
            if not defn:
                continue

            # Extract primary key constraint
            if "PRIMARY KEY" in defn.upper():
                pk_match = re.search(r"PRIMARY\s+KEY\s*\(([^)]+)\)", defn, re.IGNORECASE)
                if pk_match:
                    primary_keys.extend(
                        [col.strip() for col in pk_match.group(1).split(',')]
                    )
                continue

            # Extract foreign key constraint
            if "FOREIGN KEY" in defn.upper():
                fk_match = re.search(
                    r"FOREIGN\s+KEY\s*\(([^)]+)\)\s+REFERENCES\s+([^\s(]+)\s*\(([^)]+)\)",
                    defn,
                    re.IGNORECASE
                )
                if fk_match:
                    foreign_keys.append({
                        'columns': [col.strip() for col in fk_match.group(1).split(',')],
                        'reference_table': fk_match.group(2).strip('`"[]'),
                        'reference_columns': [col.strip() for col in fk_match.group(3).split(',')]
                    })
                continue

            # Parse regular column definition
            col_match = re.search(r"(\S+)\s+([^,]+)", defn)
            if col_match:
                name = col_match.group(1).strip('`"[]')
                definition = col_match.group(2).strip()
                
                # Skip if this is part of a constraint
                if name.upper() in ['PRIMARY', 'FOREIGN', 'CONSTRAINT']:
                    continue

                # Add inline primary key to primary_keys list
                if 'PRIMARY KEY' in definition.upper():
                    primary_keys.append(name)

                columns.append({
                    'name': name,
                    'definition': definition,
                    'nullable': 'NOT NULL' not in definition.upper()
                })

        return TableStructure(
            name=table_name,
            columns=columns,
            primary_keys=primary_keys,
            foreign_keys=foreign_keys
        )

    def _format_table_structure(self, table: TableStructure) -> str:
        """Format table structure documentation."""
        md = []
        
        # Column information table
        md.append("\n| Column | Type | Required | Key | Description |")
        md.append("|--------|------|----------|-----|-------------|")
        
        for col in table.columns:
            name = col['name']
            
            # Extract base type
            type_match = re.match(r'^(\w+(?:\s*\([^)]+\))?)', col['definition'])
            type_str = type_match.group(1) if type_match else col['definition']
            
            # Required indicator
            required = "✓" if not col['nullable'] else ""
            
            # Key information
            keys = []
            if name in table.primary_keys:
                keys.append("🔑")
            if any(name in fk['columns'] for fk in table.foreign_keys):
                keys.append("🔗")
            key_str = " ".join(keys)
            
            # Description including constraints
            desc_parts = []
            
            # Default value
            default_match = re.search(r'DEFAULT\s+([^,\s]+)', col['definition'], re.IGNORECASE)
            if default_match:
                desc_parts.append(f"Default: {default_match.group(1)}")
            
            # Check constraint
            check_match = re.search(r'CHECK\s*\(([^)]+)\)', col['definition'], re.IGNORECASE)
            if check_match:
                desc_parts.append(f"Check: {check_match.group(1)}")
            
            desc_str = ", ".join(desc_parts)
            
            md.append(f"| {name} | {type_str} | {required} | {key_str} | {desc_str} |")
        
        # Foreign key relationships
        if table.foreign_keys:
            md.append("\n#### Foreign Key Relationships\n")
            for fk in table.foreign_keys:
                md.append(f"* {', '.join(fk['columns'])} → "
                         f"{fk['reference_table']} "
                         f"({', '.join(fk['reference_columns'])})")

        return "\n".join(md)


class FileConversionException(BaseException):
    pass


class UnsupportedFormatException(BaseException):
    pass


class MarkItDown:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
        mlm_model: Optional[Any] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client
        self._mlm_model = mlm_model

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(BingSerpConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(PdfConverter())
        # 
        self.register_page_converter(EnhancedCsvConverter())
        self.register_page_converter(JsonConverter())
        self.register_page_converter(SqlConverter())
        self.register_page_converter(XmlConverter())
        self.register_page_converter(YamlConverter())

    def convert(
        self, source: Union[str, requests.Response], **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if (
                source.startswith("http://")
                or source.startswith("https://")
                or source.startswith("file://")
            ):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(
        self, path: str, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)

        for g in self._guess_ext_magic(path):
            self._append_ext(extensions, g)

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(
        self, stream: Any, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(
        self, url: str, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        response = self._requests_session.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            for g in self._guess_ext_magic(temp_path):
                self._append_ext(extensions, g)

            # Convert
            result = self._convert(temp_path, extensions, url=response.url)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(
        self, local_path: str, extensions: List[Union[str, None]], **kwargs
    ) -> DocumentConverterResult:
        error_trace = ""
        for ext in extensions + [None]:  # Try last with no extension
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)

                # Overwrite file_extension appropriately
                if ext is None:
                    if "file_extension" in _kwargs:
                        del _kwargs["file_extension"]
                else:
                    _kwargs.update({"file_extension": ext})

                # Copy any additional global options
                if "mlm_client" not in _kwargs and self._mlm_client is not None:
                    _kwargs["mlm_client"] = self._mlm_client

                if "mlm_model" not in _kwargs and self._mlm_model is not None:
                    _kwargs["mlm_model"] = self._mlm_model

                # If we hit an error log it and keep trying
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join(
                        [line.rstrip() for line in re.split(r"\r?\n", res.text_content)]
                    )
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            extensions = list()
            for g in guesses:
                ext = g.extension.strip()
                if len(ext) > 0:
                    if not ext.startswith("."):
                        ext = "." + ext
                    if ext not in extensions:
                        extensions.append(ext)
            return extensions
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return []

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
